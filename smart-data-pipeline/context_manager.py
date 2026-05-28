from __future__ import annotations

from pathlib import Path


class ContextManager:
    """Ingests context docs and exposes a concatenated business context string."""

    def __init__(self) -> None:
        self._context_chunks: dict[str, str] = {}

    def ingest_pdf(self, file_path: str) -> None:
        text = self._extract_pdf(file_path)
        self._context_chunks[Path(file_path).name] = text

    def ingest_powerpoint(self, file_path: str) -> None:
        try:
            from pptx import Presentation
        except Exception as exc:
            raise ImportError("python-pptx is required for PowerPoint ingestion") from exc

        prs = Presentation(file_path)
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        self._context_chunks[Path(file_path).name] = "\n".join(texts)

    def ingest_word(self, file_path: str) -> None:
        try:
            from docx import Document
        except Exception as exc:
            raise ImportError("python-docx is required for Word ingestion") from exc

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        self._context_chunks[Path(file_path).name] = "\n".join(paragraphs)

    def add_raw_context(self, name: str, text: str) -> None:
        self._context_chunks[name] = text.strip()

    def get_all_context(self) -> str:
        parts = []
        for name, text in self._context_chunks.items():
            parts.append(f"### SOURCE: {name}\n{text}")
        return "\n\n".join(parts)

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        try:
            import pdfplumber

            content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    content.append(page.extract_text() or "")
            return "\n".join(content)
        except Exception as exc:
            raise RuntimeError(f"Unable to extract PDF text: {exc}") from exc
